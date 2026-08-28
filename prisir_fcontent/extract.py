# -*- coding: utf-8 -*-
"""内容抽取:按扩展名分派,把文件正文抽成纯文本。

支持:
- 纯文本/代码/标记: .txt .md .log .py .js .ts .json .html .css .xml .csv 等直读(截断)。
- .docx  python-docx
- .pdf   pypdf
- .pptx  python-pptx
- 图片(.png/.jpg/.jpeg/.bmp/.webp): OCR(rapidocr_onnxruntime),可选开关,默认关。

红线:
- 每文件内容截断 MAX_CHARS(防大文件爆库/爆内存)。
- 抽不出文本的扩展名(音视频/可执行/压缩包)直接跳过不索引。
- 加密/损坏/解析失败返 None(跳过,不阻塞批量),绝不抛到调用方。
- 图片 OCR 默认关:只有装了 rapidocr 且 enable(ocr=True) 显式开才走 OCR。
- OCR 置信门控:行 score < OCR_MIN_SCORE 的行不入库(宁可漏不可错)。
"""
import os

MAX_CHARS = 512 * 1024  # 512KB 字符截断

# OCR 置信门控:RapidOCR 行 score < 此值的行丢弃(错字比漏字更伤搜索信任)
OCR_MIN_SCORE = 0.8

# 图片类:OCR 抽取(可选,默认关)
_IMG_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".webp"}

# 纯文本/代码/标记类:直接按文本读(utf-8,errors=replace)
_TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".log", ".rst",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".h", ".cpp", ".hpp",
    ".go", ".rs", ".rb", ".php", ".sh", ".bat", ".ps1",
    ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".html", ".htm", ".css", ".xml", ".svg",
    ".csv", ".tsv",
}

# 二进制可解析类:走对应解析库
_DOCX_EXT = ".docx"
_PDF_EXT = ".pdf"
_PPTX_EXT = ".pptx"


def _read_text_file(path: str) -> str:
    with open(path, encoding="utf-8", errors="replace") as f:
        return f.read(MAX_CHARS)


def _extract_docx(path: str) -> str:
    import docx  # python-docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs]
    for tbl in d.tables:
        for row in tbl.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return "\n".join(parts)[:MAX_CHARS]


def _extract_pdf(path: str) -> str:
    from pypdf import PdfReader
    r = PdfReader(path)
    parts = []
    for pg in r.pages:
        try:
            parts.append(pg.extract_text() or "")
        except Exception:  # noqa: BLE001
            continue  # 单页失败不阻塞整篇
    return "\n".join(parts)[:MAX_CHARS]


def _extract_pptx(path: str) -> str:
    from pptx import Presentation  # python-pptx
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    parts.append("".join(run.text for run in para.runs))
    return "\n".join(parts)[:MAX_CHARS]


# ---- 图片 OCR(可选,默认关) ----
_OCR = None          # RapidOCR 单例(lazy,首载 ~3s)
_OCR_TRIED = False   # 是否已尝试初始化(避免每次 import 失败重试)


def ocr_available():
    """探测 OCR 能力:rapidocr_onnxruntime 是否可 import。返 (bool, reason)。"""
    try:
        import rapidocr_onnxruntime  # noqa: F401
        return True, "ok"
    except Exception:  # noqa: BLE001
        return False, "module_not_installed"


def _get_ocr():
    """lazy 初始化 RapidOCR 单例;不可用返 None。首载较重(~3s),复用不重建。"""
    global _OCR, _OCR_TRIED
    if _OCR is not None or _OCR_TRIED:
        return _OCR
    _OCR_TRIED = True
    ok, _ = ocr_available()
    if not ok:
        return None
    try:
        from rapidocr_onnxruntime import RapidOCR
        _OCR = RapidOCR()
    except Exception:  # noqa: BLE001
        _OCR = None
    return _OCR


def _extract_ocr(path: str) -> str:
    """图片 OCR 抽取。行 score<OCR_MIN_SCORE 的门控丢弃;无文本/失败返 None。"""
    ocr = _get_ocr()
    if ocr is None:
        return None
    res, _ = ocr(path)
    if not res:
        return None
    parts = [r[1] for r in res if len(r) >= 3 and r[2] >= OCR_MIN_SCORE]
    if not parts:
        return None
    return "\n".join(parts)[:MAX_CHARS]


def extract_text(path: str, ocr: bool = False):
    """抽取文件正文为纯文本。不支持的扩展名/失败返 None(跳过)。

    ocr=True 且能力可用时,图片(.png/.jpg/...)走 OCR;默认 False 图片跳过。"""
    if not path or not isinstance(path, str):
        return None
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext in _TEXT_EXTS:
            return _read_text_file(path)
        if ext == _DOCX_EXT:
            return _extract_docx(path)
        if ext == _PDF_EXT:
            return _extract_pdf(path)
        if ext == _PPTX_EXT:
            return _extract_pptx(path)
        if ext in _IMG_EXTS:
            return _extract_ocr(path) if ocr else None
        return None  # 不支持的类型(音视频/可执行等)跳过
    except Exception:  # noqa: BLE001
        return None  # 加密/损坏/解析失败:跳过不阻塞


def supported(ext: str, ocr: bool = False) -> bool:
    """该扩展名是否支持内容抽取(供扫描时预过滤)。图片仅在 ocr=True 时纳入。"""
    e = (ext or "").lower()
    if e in _TEXT_EXTS or e in {_DOCX_EXT, _PDF_EXT, _PPTX_EXT}:
        return True
    if e in _IMG_EXTS:
        return ocr and ocr_available()[0]
    return False
